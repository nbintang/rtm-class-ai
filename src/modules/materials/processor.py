from pathlib import Path

from src.core.exceptions import AppError


def sanitize_text(text: str) -> str:
    return " ".join(text.split()).strip()


def extract_pdf_text(file_path: str) -> str:
    try:
        from pypdf import PdfReader
    except Exception as exc:  # pragma: no cover
        raise AppError(f"pypdf is required for PDF processing: {exc}") from exc

    path = Path(file_path)
    if not path.exists():
        raise AppError(f"PDF file not found: {file_path}", status_code=404, code="file_not_found")

    reader = PdfReader(str(path))
    content = " ".join(page.extract_text() or "" for page in reader.pages)
    return sanitize_text(content)


def extract_ppt_text(file_path: str) -> str:
    try:
        from pptx import Presentation
    except Exception as exc:  # pragma: no cover
        raise AppError(f"python-pptx is required for PPT processing: {exc}") from exc

    path = Path(file_path)
    if not path.exists():
        raise AppError(f"PPT file not found: {file_path}", status_code=404, code="file_not_found")

    prs = Presentation(str(path))
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text:
                parts.append(text)
    return sanitize_text(" ".join(parts))


def extract_text_from_content(content: str) -> str:
    cleaned = sanitize_text(content)
    if not cleaned:
        raise AppError("Material content is empty after cleaning", status_code=400, code="empty_content")
    return cleaned

