from __future__ import annotations

from io import BytesIO
from pathlib import Path


def _normalize_text(text: str) -> str:
    return " ".join(text.replace("\x00", " ").split())


def extract_text_from_pdf(payload: bytes) -> str:
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise ValueError("PDF support is unavailable. Install pypdf.") from exc

    try:
        reader = PdfReader(BytesIO(payload))
    except Exception as exc:
        raise ValueError(f"Failed to read PDF file: {exc}") from exc

    chunks: list[str] = []
    for page in reader.pages:
        text = page.extract_text() or ""
        if text.strip():
            chunks.append(text)
    return _normalize_text("\n".join(chunks))


def extract_text_from_pptx(payload: bytes) -> str:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise ValueError("PPTX support is unavailable. Install python-pptx.") from exc

    try:
        presentation = Presentation(BytesIO(payload))
    except Exception as exc:
        raise ValueError(f"Failed to read PPTX file: {exc}") from exc

    chunks: list[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text and text.strip():
                chunks.append(text)
    return _normalize_text("\n".join(chunks))


def extract_text_from_txt(payload: bytes) -> str:
    try:
        text = payload.decode("utf-8-sig")
    except UnicodeDecodeError:
        text = payload.decode("utf-8", errors="replace")
    return _normalize_text(text)


def extract_material_text(
    *,
    filename: str,
    content_type: str | None,
    payload: bytes,
) -> tuple[str, str, list[str]]:
    warnings: list[str] = []
    ext = Path(filename).suffix.lower()

    if ext == ".pdf":
        file_type = "pdf"
        text = extract_text_from_pdf(payload)
    elif ext == ".pptx":
        file_type = "pptx"
        text = extract_text_from_pptx(payload)
    elif ext == ".txt":
        file_type = "txt"
        text = extract_text_from_txt(payload)
    else:
        raise ValueError(
            "Unsupported file type. Allowed extensions: .pdf, .pptx, .txt"
        )

    if not text.strip():
        raise ValueError("Extracted text is empty.")

    if content_type and file_type == "txt" and "text" not in content_type.lower():
        warnings.append("Uploaded file extension is .txt but content-type is unusual.")

    return text, file_type, warnings
